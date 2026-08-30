#!/usr/bin/env python3
"""Generic GAIA seed-skill evaluator (single-agent).

Relative to the plain rollout it adds:
- ``--input``  any JSONL of GAIA items (e.g. ``test65.jsonl``)
- ``--model``  any OpenAI/OpenRouter chat model (e.g. ``openai/gpt-4o``)
- ``--workers`` parallel inference threads
- attachment handling: text inline, xlsx→markdown, pdf→per-page text,
  docx/pptx→text, png/jpg→image_url for vision models, audio→fallback note

Usage:
    OPENAI_API_KEY=<key> python -m maskills.envs.gaia.single_agent \
        --input data/gaia/test65.jsonl \
        --model openai/gpt-4o \
        --workers 6 \
        --out data/gaia/test65_gpt4o.jsonl
"""

from __future__ import annotations

import argparse
import base64
import json
import os
import re
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from openai import OpenAI

from ._keys import require_api_keys

#: Where GAIA attachments are resolved from, relative to a source checkout.
#: Skill libraries are never defaulted: every entry point takes the directory,
#: because no library ships with the repository.
REPO_ROOT = Path(__file__).resolve().parents[3]
DATA_DIR = REPO_ROOT / "data" / "gaia"

# ── Skill loading ──────────────────────────────────────────────────────────

FRONTMATTER_RE = re.compile(r"^---\n(.*?)\n---\n(.*)$", re.DOTALL)


def parse_skill(path: Path) -> dict:
    text = path.read_text()
    m = FRONTMATTER_RE.match(text)
    if not m:
        raise ValueError(f"no frontmatter in {path}")
    fm, body = m.group(1), m.group(2)
    name = re.search(r"^name:\s*(.+)$", fm, re.M).group(1).strip()
    desc = re.search(r"^description:\s*(.+)$", fm, re.M).group(1).strip()
    return {"name": name, "description": desc, "body": body, "path": str(path)}


def load_skills(skills_dir: Path) -> dict:
    """Read ``<skills_dir>/SKILL.md`` plus one sub-skill per sub-directory."""
    skills_dir = Path(skills_dir)
    if not (skills_dir / "SKILL.md").exists():
        raise SystemExit(f"No root SKILL.md under {skills_dir}.")
    root = parse_skill(skills_dir / "SKILL.md")
    skills = {}
    for sub in sorted(skills_dir.iterdir()):
        if sub.is_dir():
            sk = parse_skill(sub / "SKILL.md")
            skills[sub.name] = sk
    return {"root": root, "skills": skills}


# ── Attachment handlers ────────────────────────────────────────────────────


def _read_text(p: Path) -> str:
    return p.read_text(errors="ignore")


def _read_xlsx_as_markdown(p: Path, max_rows: int = 200) -> str:
    """Render every sheet of an xlsx as a markdown table."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "(openpyxl not available; cannot parse xlsx)"
    wb = load_workbook(p, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"\n### Sheet: `{sheet.title}` ({sheet.max_row} rows × {sheet.max_column} cols)")
        rows = list(sheet.iter_rows(values_only=True))
        rows = rows[:max_rows]
        if not rows:
            parts.append("(empty)")
            continue
        ncols = max(len(r) for r in rows)
        header = [str(c) if c is not None else "" for c in rows[0]] + [""] * (ncols - len(rows[0]))
        parts.append("| " + " | ".join(header) + " |")
        parts.append("|" + "|".join("---" for _ in range(ncols)) + "|")
        for r in rows[1:]:
            cells = [str(c) if c is not None else "" for c in r] + [""] * (ncols - len(r))
            parts.append("| " + " | ".join(cells) + " |")
        if len(list(sheet.iter_rows())) > max_rows:
            parts.append(f"(... truncated at {max_rows} rows)")
    return "\n".join(parts)


def _read_pdf_as_text(p: Path, max_pages: int = 30) -> str:
    """Per-page text dump using pdfplumber."""
    try:
        import pdfplumber
    except ImportError:
        return "(pdfplumber not available; cannot parse pdf)"
    parts = []
    with pdfplumber.open(p) as pdf:
        total = len(pdf.pages)
        for i, page in enumerate(pdf.pages[:max_pages]):
            txt = page.extract_text() or ""
            parts.append(f"\n--- Page {i+1}/{total} ---\n{txt}")
        if total > max_pages:
            parts.append(f"\n(... truncated at page {max_pages} of {total})")
    return "\n".join(parts)


def _read_docx_as_text(p: Path) -> str:
    try:
        import docx
    except ImportError:
        return "(python-docx not available; cannot parse docx)"
    d = docx.Document(p)
    parts = [para.text for para in d.paragraphs]
    for ti, t in enumerate(d.tables):
        parts.append(f"\n### Table {ti+1}")
        for row in t.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _read_pptx_as_text(p: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "(python-pptx not available; cannot parse pptx)"
    pres = Presentation(p)
    parts = []
    for i, slide in enumerate(pres.slides, 1):
        parts.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _image_as_data_url(p: Path) -> str:
    ext = p.suffix.lower().lstrip(".")
    mime = {"jpg": "jpeg", "jpeg": "jpeg", "png": "png"}.get(ext, "png")
    b64 = base64.b64encode(p.read_bytes()).decode("ascii")
    return f"data:image/{mime};base64,{b64}"


# ── Prompt building ────────────────────────────────────────────────────────

EXT_TO_SKILLS = {
    ".png": ["image_vision"],
    ".jpg": ["image_vision"],
    ".jpeg": ["image_vision"],
    ".xlsx": ["table_xlsx_csv"],
    ".csv": ["table_xlsx_csv"],
    ".tsv": ["table_xlsx_csv"],
    ".pdf": ["pdf_extract"],
    ".docx": ["pdf_extract"],   # closest analogue
    ".pptx": ["pdf_extract"],
    ".py": ["code_python"],
    ".txt": [],
    ".json": [],
    ".jsonld": [],
    ".xml": [],
    ".pdb": [],
    ".zip": [],
    ".mp3": [],
    ".m4a": [],
    ".mov": [],
}

TIER_A = [
    "tool_call_syntax",        # B1 — read first; syntax discipline
    "attachment_first",        # B2 — pre-inlined attachments, no fake file-open tools
    "constraint_checklist",
    "plan_and_decompose",
    "cross_verify_fact",
    "answer_format",
    "verify_before_finalize",
    "critique_and_revise",
]


def render_system_prompt(skill_bundle: dict, relevant_b: list) -> str:
    root = skill_bundle["root"]
    skills = skill_bundle["skills"]
    parts = [root["body"].strip(), "\n\n---\n\n# Tier A skill bodies (always-on)\n"]
    for n in TIER_A:
        sk = skills[n]
        parts.append(f"\n## skill: `{sk['name']}`\n{sk['body'].strip()}\n")
    parts.append("\n---\n\n# Tier B skill bodies (relevant for this task)\n")
    for n in relevant_b:
        sk = skills[n]
        parts.append(f"\n## skill: `{sk['name']}`\n{sk['body'].strip()}\n")
    if not relevant_b:
        parts.append("\n(no Tier B skill loaded — pure-reasoning task)\n")
    not_loaded = [n for n in skills if n not in TIER_A and n not in relevant_b]
    if not_loaded:
        parts.append(
            "\n---\n\n# Tier B catalog (NOT loaded — metadata only)\n\n"
            + "\n".join(
                f"- `{skills[n]['name']}` — {skills[n]['description']}"
                for n in not_loaded
            )
        )
    return "".join(parts)


TEXT_EXTS = {".txt", ".csv", ".tsv", ".py", ".json", ".jsonld", ".xml", ".pdb"}
TABLE_EXTS = {".xlsx"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg"}
UNSUPPORTED_EXTS = {".mp3", ".m4a", ".mov", ".zip"}


def render_user_message(item: dict) -> list:
    """Returns a list of OpenAI content parts (text + optional image_url).

    For text-only attachments / no attachment, a single text part is returned;
    for image attachments, a text part followed by an image_url part.
    """
    parts_text = [f"# GAIA question (task_id = {item['task_id']}, Level = {item['Level']})\n"]
    parts_text.append(item["Question"].strip())
    image_part = None

    if item.get("file_name"):
        attach_path = DATA_DIR / item["file_path"]
        ext = os.path.splitext(item["file_name"])[1].lower()
        parts_text.append(f"\n\n---\n\n## Attached file: `{item['file_name']}`\n")
        if ext in TEXT_EXTS:
            try:
                content = _read_text(attach_path)
                parts_text.append(f"```{ext.lstrip('.')}\n{content}\n```")
            except Exception as e:  # noqa: BLE001
                parts_text.append(f"(could not read: {e})")
        elif ext in TABLE_EXTS:
            try:
                content = _read_xlsx_as_markdown(attach_path)
                parts_text.append(f"(parsed by openpyxl, all sheets as markdown)\n{content}")
            except Exception as e:  # noqa: BLE001
                parts_text.append(f"(xlsx parse error: {e})")
        elif ext in PDF_EXTS:
            try:
                content = _read_pdf_as_text(attach_path)
                parts_text.append(f"(parsed by pdfplumber, per-page text)\n{content}")
            except Exception as e:  # noqa: BLE001
                parts_text.append(f"(pdf parse error: {e})")
        elif ext in DOCX_EXTS:
            try:
                content = _read_docx_as_text(attach_path)
                parts_text.append(f"(parsed by python-docx)\n{content}")
            except Exception as e:  # noqa: BLE001
                parts_text.append(f"(docx parse error: {e})")
        elif ext in PPTX_EXTS:
            try:
                content = _read_pptx_as_text(attach_path)
                parts_text.append(f"(parsed by python-pptx)\n{content}")
            except Exception as e:  # noqa: BLE001
                parts_text.append(f"(pptx parse error: {e})")
        elif ext in IMAGE_EXTS:
            parts_text.append("(image attached as multimodal input below)")
            image_part = {
                "type": "image_url",
                "image_url": {"url": _image_as_data_url(attach_path)},
            }
        elif ext in UNSUPPORTED_EXTS:
            parts_text.append(
                f"(Attachment type `{ext}` not supported in this text-only "
                f"harness — no transcription/extraction available. State this "
                f"and best-guess from question text alone.)"
            )
        else:
            parts_text.append(f"(Unknown extension `{ext}`; cannot parse.)")

    parts_text.append(
        "\n\nWork through the problem following the protocol in the root "
        "SKILL.md. End with a single `FINAL ANSWER:` line."
    )
    full_text = "".join(parts_text)
    if image_part is None:
        return [{"type": "text", "text": full_text}]
    else:
        return [{"type": "text", "text": full_text}, image_part]


def select_relevant_b(item: dict) -> list:
    rel = []
    fn = item.get("file_name", "")
    if fn:
        ext = os.path.splitext(fn)[1].lower()
        rel.extend(EXT_TO_SKILLS.get(ext, []))
    q = item["Question"].lower()
    if any(k in q for k in ("how many", "what percent", "round", "calculate", "percentage")):
        if "calc_and_units" not in rel:
            rel.append("calc_and_units")
    if any(k in q for k in (
        "wikipedia", "according to", "paper", "github", "arxiv", "youtube",
        "video", "website", "url", "page",
    )):
        for n in ("web_search", "web_browse_deep"):
            if n not in rel:
                rel.append(n)
    return rel


# ── Answer parsing & scoring ───────────────────────────────────────────────


def parse_final_answer(text: str) -> str:
    matches = list(re.finditer(r"FINAL ANSWER:\s*(.+?)\s*$", text, re.M))
    if not matches:
        return ""
    return matches[-1].group(1).strip().rstrip(".").strip('"').strip()


def normalize_for_match(s: str) -> str:
    s = s.strip().rstrip(".").strip('"').strip("'").strip()
    s_no_comma = s.replace(",", "")
    try:
        f = float(s_no_comma)
        return f"{f:g}"
    except ValueError:
        return s.lower()


def is_correct(pred: str, gold: str) -> bool:
    return normalize_for_match(pred) == normalize_for_match(gold)


# ── Worker ─────────────────────────────────────────────────────────────────


def run_one(client: OpenAI, item: dict, bundle: dict, model: str,
            max_tokens: int) -> dict:
    rel_b = select_relevant_b(item)
    sys_prompt = render_system_prompt(bundle, rel_b)
    usr_parts = render_user_message(item)
    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": usr_parts},
            ],
            temperature=0.0,
            max_tokens=max_tokens,
        )
        content = resp.choices[0].message.content or ""
        usage = resp.usage
        in_tok = usage.prompt_tokens if usage else 0
        out_tok = usage.completion_tokens if usage else 0
        err = ""
    except Exception as e:  # noqa: BLE001
        content = ""
        in_tok = out_tok = 0
        err = f"{type(e).__name__}: {e}"

    pred = parse_final_answer(content)
    ok = is_correct(pred, item["Final answer"]) if pred else False
    return {
        "task_id": item["task_id"],
        "Level": item["Level"],
        "kind": item["kind"],
        "rel_B": rel_b,
        "gold": item["Final answer"],
        "pred": pred,
        "correct": ok,
        "in_tok": in_tok,
        "out_tok": out_tok,
        "error": err,
        "raw_response": content,
    }


# ── Main ───────────────────────────────────────────────────────────────────


def main():
    require_api_keys()
    ap = argparse.ArgumentParser()
    ap.add_argument("--skills-dir", type=str, required=True,
                    help="skill library to run: a directory with a root "
                         "SKILL.md and one sub-directory per sub-skill")
    ap.add_argument("--input", type=str, required=True,
                    help="JSONL of GAIA items")
    ap.add_argument("--model", type=str, default="openai/gpt-4o",
                    help="OpenRouter chat model ID")
    ap.add_argument("--out", type=str, default="",
                    help="results JSONL (default: derived from --input)")
    ap.add_argument("--workers", type=int, default=6)
    ap.add_argument("--n", type=int, default=0, help="0=all; else limit")
    ap.add_argument("--max-tokens", type=int, default=3000)
    args = ap.parse_args()

    bundle = load_skills(Path(args.skills_dir))
    print(f"Skills loaded: 1 root + {len(bundle['skills'])} sub-skills "
          f"({', '.join(bundle['skills'])})")

    items = [json.loads(l) for l in Path(args.input).open()]
    if args.n > 0:
        items = items[: args.n]
    out_path = Path(args.out) if args.out else (
        Path(args.input).with_suffix("")
        .with_suffix(f".{args.model.replace('/', '_')}.jsonl")
    )

    print(f"Input    : {args.input} ({len(items)} items)")
    print(f"Model    : {args.model}")
    print(f"Workers  : {args.workers}")
    print(f"Output   : {out_path}\n")

    client = OpenAI()
    results = [None] * len(items)
    t0 = time.time()

    with ThreadPoolExecutor(max_workers=args.workers) as ex:
        fut_to_idx = {
            ex.submit(run_one, client, item, bundle, args.model, args.max_tokens): i
            for i, item in enumerate(items)
        }
        done = 0
        for fut in as_completed(fut_to_idx):
            i = fut_to_idx[fut]
            r = fut.result()
            results[i] = r
            done += 1
            mark = "✓" if r["correct"] else ("E" if r["error"] else "✗")
            print(f"  [{done:3d}/{len(items)}] {mark} L{r['Level']} "
                  f"{r['task_id'][:8]} {r['kind']:6s}  pred={r['pred']!r:25s}"
                  f"  gold={r['gold']!r:25s}  ({r['in_tok']}+{r['out_tok']} tok)"
                  + (f"  err={r['error']}" if r["error"] else ""))

    out_path.write_text("\n".join(json.dumps(r, ensure_ascii=False) for r in results))
    dt = time.time() - t0

    correct = sum(1 for r in results if r["correct"])
    total_in = sum(r["in_tok"] for r in results)
    total_out = sum(r["out_tok"] for r in results)
    errors = sum(1 for r in results if r["error"])

    print("\n" + "=" * 78)
    print(f"OVERALL  {correct}/{len(items)} = {correct/len(items)*100:5.1f}%"
          f"   ({errors} API errors)   in {dt:.1f}s")
    print(f"Tokens   {total_in:,} in + {total_out:,} out = {total_in+total_out:,}")

    by_lvl = {}
    for r in results:
        by_lvl.setdefault(r["Level"], [0, 0])
        by_lvl[r["Level"]][0] += int(r["correct"])
        by_lvl[r["Level"]][1] += 1
    print("\nBy Level:")
    for lvl in sorted(by_lvl):
        c, t = by_lvl[lvl]
        print(f"  L{lvl}: {c:2d}/{t:2d}  ({c/t*100:5.1f}%)")

    by_kind = {}
    for r in results:
        by_kind.setdefault(r["kind"], [0, 0])
        by_kind[r["kind"]][0] += int(r["correct"])
        by_kind[r["kind"]][1] += 1
    print("\nBy modality:")
    for k in sorted(by_kind):
        c, t = by_kind[k]
        print(f"  {k:8s}: {c:2d}/{t:2d}  ({c/t*100:5.1f}%)")
    print(f"\nResults → {out_path}")


if __name__ == "__main__":
    main()
